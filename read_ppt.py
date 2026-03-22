import codecs
from pptx import Presentation

try:
    prs = Presentation('Portfolio Creation Lecture 2.pptx')
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_runs.append(shape.text)
    
    with codecs.open('ppt_text.txt', 'w', 'utf-8') as f:
        f.write('\n'.join(text_runs))
    print("Successfully extracted text")
except Exception as e:
    print(f"Error: {e}")
